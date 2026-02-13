# backend/app/services/notifications_pptx.py
from typing import List, Dict, Any, Optional
import os
from datetime import datetime
from pptx import Presentation
from app.config import settings

def _slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            t = shape.text.strip()
            if t:
                texts.append(t)
    return "\n".join(texts).strip()

def load_notifications_from_pptx(query: Optional[str] = None) -> List[Dict[str, Any]]:
    path = settings.NOTIFICATIONS_PPTX_PATH
    if not os.path.exists(path):
        return []

    prs = Presentation(path)
    mtime = datetime.fromtimestamp(os.path.getmtime(path))
    updated_at_iso = mtime.isoformat()

    items: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide)
        if not text:
            continue

        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        title = lines[0] if lines else f"Slide {idx}"
        preview = text[:200] + ("..." if len(text) > 200 else "")

        items.append({
            "id": idx,
            "title": title,
            "text": preview,
            "content": text,
            "date": updated_at_iso,
        })

    if query:
        q = query.strip().lower()
        items = [x for x in items if q in (x["title"] or "").lower() or q in (x["content"] or "").lower()]

    return items
